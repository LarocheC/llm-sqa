"""
Build the Plan B results deck in two formats from ONE shared content model:
  - planb_deck.pptx  (native python-pptx: editable text boxes, tables, embedded figures)
  - planb_deck.pdf   (matplotlib PdfPages: vector text + raster figures)

Run with the deck venv:
  experiments/.deckenv/bin/python experiments/build_deck.py
"""
import os, textwrap

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import matplotlib.image as mpimg

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RES = os.path.join(os.path.dirname(os.path.abspath(__file__)), "results")
PLB = os.path.join(RES, "planb")

# ---- palette ----
BG   = "#0d1117"; PANEL = "#161b22"; INK = "#e6edf3"; DIM = "#9aa7b4"
LINE = "#30363d"; ACC = "#58a6ff"; GOOD = "#3fb950"; BAD = "#f85149"
WARN = "#d29922"; MAG = "#bc8cff"
def H(c): return RGBColor.from_string(c.lstrip("#").upper())
KIND_COLOR = {"h3":MAG,"bullet":INK,"sub":DIM,"plain":INK,"good":GOOD,
              "bad":BAD,"warn":WARN,"dim":DIM,"acc":ACC}

# ======================================================================
# diagrams rendered once as PNGs (embedded by both backends)
# ======================================================================
def _box(ax, x, y, w, h, text, sub, edge, fs=12):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.04",
                                linewidth=2, edgecolor=edge, facecolor=PANEL))
    ax.text(x+w/2, y+h*0.62, text, ha="center", va="center", color=INK, fontsize=fs, weight="bold")
    if sub:
        ax.text(x+w/2, y+h*0.27, sub, ha="center", va="center", color=DIM, fontsize=fs-3)

def _arrow(ax, x1, y1, x2, y2):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>", mutation_scale=14,
                                 lw=2, color=ACC))

def make_arch():
    fig, ax = plt.subplots(figsize=(6.2, 5.4), facecolor=BG); ax.set_facecolor(BG)
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")
    _box(ax, 0.3, 8.6, 4.3, 1.1, "Whisper-Large-v2", "frozen · mel→1280-d", WARN)
    _box(ax, 5.4, 8.6, 4.3, 1.1, "BEATs", "frozen · wav→768-d", WARN)
    _box(ax, 2.8, 6.7, 4.4, 0.95, "concat → 2048-d", "", LINE)
    _box(ax, 2.4, 4.8, 5.2, 1.1, "Window Q-Former", "TRAINED · ~88 tokens (17:1)", GOOD)
    _box(ax, 2.7, 3.1, 4.6, 0.95, "Linear proj 768→4096", "TRAINED", GOOD)
    _box(ax, 1.5, 1.0, 7.0, 1.25, "Vicuna-7B  (FROZEN)", "", WARN)
    ax.text(5.0, 1.35, "+ LoRA r=8 α=28 · TRAINED", ha="center", va="center", color=GOOD, fontsize=10)
    _arrow(ax, 2.4, 8.6, 4.0, 7.7); _arrow(ax, 7.6, 8.6, 6.0, 7.7)
    _arrow(ax, 5.0, 6.7, 5.0, 5.95); _arrow(ax, 5.0, 4.8, 5.0, 4.1); _arrow(ax, 5.0, 3.1, 5.0, 2.3)
    fig.tight_layout(pad=0.2)
    p = os.path.join(PLB, "_diag_arch.png"); fig.savefig(p, dpi=200, facecolor=BG); plt.close(fig); return p

def make_pipeline():
    fig, ax = plt.subplots(figsize=(12.4, 2.2), facecolor=BG); ax.set_facecolor(BG)
    ax.set_xlim(0, 40); ax.set_ylim(0, 6); ax.axis("off")
    stages = [("Stage 0 — probe", "decodable from\nfrozen features?", MAG),
              ("Corpus generation", "on-the-fly degrad.\n+ Opus paraphrase", ACC),
              ("Two-stage LoRA-SFT", "calibrate → reason\nencoders frozen", GOOD),
              ("Held-out eval", "re-measure all\n5 findings", WARN)]
    x = 0.4; w = 8.6; gap = 1.7
    cx = []
    for t, s, e in stages:
        ax.add_patch(FancyBboxPatch((x, 1.4), w, 3.0, boxstyle="round,pad=0.1,rounding_size=0.3",
                                    lw=2, edgecolor=e, facecolor=PANEL))
        ax.text(x+w/2, 3.5, t, ha="center", color=INK, fontsize=12, weight="bold")
        ax.text(x+w/2, 2.3, s, ha="center", color=DIM, fontsize=9.5)
        cx.append((x, x+w)); x += w + gap
    for i in range(len(stages)-1):
        _arrow(ax, cx[i][1], 2.9, cx[i+1][0], 2.9)
    fig.tight_layout(pad=0.1)
    p = os.path.join(PLB, "_diag_pipe.png"); fig.savefig(p, dpi=200, facecolor=BG); plt.close(fig); return p

ARCH = make_arch(); PIPE = make_pipeline()

# ======================================================================
# shared content model
# B(kind, text): body line.  Slides may carry image / table / kpis / cols.
# ======================================================================
def B(k, t): return (k, t)

SLIDES = [
 # 1 title
 {"type":"title",
  "title":"Plan B",
  "subtitle":"De-specializing SALMONN-SQA into a calibrated, multi-dimensional speech-quality rater",
  "body":[B("plain","Turning a lenient noise specialist into a model that perceives, names, and scores"),
          B("plain","every common degradation — and emits a de-compressed, calibrated MOS."),
          B("dim",""),
          B("dim","Clement Laroche · final model ckpt_stage2_v3 · LibriTTS-R → VoiceBank-DEMAND held-out")]},

 # 2 the model
 {"title":"The model: SALMONN-7B, SQA-finetuned",
  "body":[B("plain","Audio LLM: two frozen encoders → a trained audio bottleneck → frozen Vicuna-7B + small LoRA."),
          B("bullet","Whisper-Large-v2 — mel→1280-d. ASR-trained, deliberately reverb-invariant."),
          B("bullet","BEATs — raw waveform → 768-d acoustic-event features."),
          B("bullet","Window Q-Former — ~0.33s windows → ~88 tokens. TRAINED. ~17:1 compression."),
          B("bullet","Linear projection 768→4096, spliced at <SpeechHere>. TRAINED."),
          B("bullet","Vicuna-7B-v1.5 + LoRA (r=8, α=28). Base FROZEN, LoRA TRAINED."),
          B("dim","Plan B trains only the green blocks + LoRA. Encoders stay frozen.")],
  "image":ARCH},

 # 3 the problem
 {"title":"The problem: a lenient noise specialist",
  "body":[B("plain","Good additive-noise descriptions — but as a quality meter, weak and blind to whole classes."),
          B("h3","① MOS isn't a usable scale"),
          B("bullet","Coarse 5 discrete values, floored at 2.5. Agrees with PESQ/NISQA/DNSMOS at only ρ 0.40–0.49"),
          B("sub","— yet those metrics agree with each other at 0.72–0.82. It is the outlier."),
          B("h3","② Blind to whole degradation types"),
          B("bullet","Reverb ρ(MOS,severity) = −0.11 (no response). Weak on bandwidth (−0.39), clipping (−0.33)."),
          B("h3","③ Blind to enhancement"),
          B("bullet","Denoise: every objective metric improves, MOS moves +0.03."),
          B("h3","④ Brittle output"),
          B("bullet","OOD degradation → made-up JSON schema. Acquiescent: leading yes/no → always YES."),
          B("good","Goal: keep the descriptive ability; add perception + calibration across ALL degradations.")]},

 # 4 evidence
 {"title":"Evidence of the blindness (original model)",
  "body":[B("bullet","Left: MOS barely tracks SNR, and sits in a narrow high band — the floored, compressed scale."),
          B("bullet","Right: vs DNSMOS / NISQA / PESQ the MOS quantizes into a few horizontal stripes (ρ 0.40–0.49).")],
  "image_pair":[os.path.join(RES,"mos_vs_snr.png"), os.path.join(RES,"mos_vs_neural.png")]},

 # 5 experiment design
 {"title":"Experiment design: de-risk, then train",
  "topimage":PIPE,
  "body":[B("h3","Ask the key question cheaply, first"),
          B("bullet","Whisper is reverb-invariant; Q-Former compresses 17:1. If reverb/bandwidth info is destroyed"),
          B("sub","in the frozen front-end, no LoRA training recovers it → would need to unfreeze encoders (heavy)."),
          B("h3","The hypothesis"),
          B("bullet","Blindness lives in the LLM's learned read-out (only ever trained on noise), not the architecture."),
          B("good","If true, the fix is data + targets — not surgery on the encoders.")]},

 # 6 stage0 method
 {"title":"Stage 0 — is degradation decodable from frozen features?",
  "body":[B("plain","Known-parameter degradations → features at 4 frozen taps → linear probe (PCA+ridge),"),
          B("plain","scored by R²+Spearman, held out by utterance."),
          B("dim","A linear probe is the point: tests whether info is linearly decodable, not buried."),
          B("h3","Taps: whisper · beats · concat · qformer (the bottleneck that reaches the LLM)"),
          B("good","Result — PASSED. R² 0.82–0.95, Spearman 0.90–0.98 at every tap, incl. Q-Former output."),
          B("bullet","Decision: keep encoders frozen; the fix is data + targets.")],
  "image":os.path.join(PLB,"stage0_r2_bars.png")},

 # 7 stage0 scatter
 {"title":"Stage 0 — the signal survives to the LLM's doorstep",
  "body":[B("h3","Where the blindness actually is"),
          B("bullet","Predicted-vs-true RT60 / cutoff / clip-fraction sit on y=x at every tap, incl. qformer."),
          B("bullet","The degradation reaches the LLM intact — the model was never taught to read it"),
          B("sub","(its read-out head was trained only on additive noise)."),
          B("dim","Each axis is probed one-at-a-time → proves per-axis detectability, not type-discrimination."),
          B("dim","Discriminating WHICH degradation is the residual that v3's synthetic-type data fixes.")],
  "image":os.path.join(PLB,"stage0_scatter.png")},

 # 8 data generation
 {"title":"Data generation: on-the-fly, known-parameter corpus",
  "body":[B("h3","Clean speech — leakage-safe"),
          B("bullet","Train: LibriTTS-R train-clean-100 — 4,000 clips. Val: dev-clean — 300 clips."),
          B("bullet","No VCTK speakers → disjoint from VoiceBank-DEMAND eval. No DEMAND noise in training."),
          B("h3","Description text — Opus 4.8 paraphrase"),
          B("bullet","~740 unique profiles → one Opus 4.8 call each → grounding-verified prose"),
          B("bullet","All 4,300 descriptions paraphrased, 0 template fallbacks.")],
  "table":{"colw":[0.20,0.62,0.18],"cols":["axis","how","train"],
           "rows":[["reverb","real measured RIRs + synth exp-decay","2,783"],
                   ["noise","real MUSAN + synth white/pink/brown","512"],
                   ["codec","real Opus/MP3 → bandwidth","509"],
                   ["bandwidth","Butterworth low-pass","DSP"],
                   ["clipping / disc.","hard-clip / frame-drop","DSP"],
                   ["loudness","re-gain","DSP"],
                   ["clean","untouched","373"]]},
  "footnote":"real + synthetic mix is the v3 addition — closes the train/eval gap (eval sweep uses synthetic types)."},

 # 9 artefacts / severity map
 {"title":"Choice of artefacts & the severity map",
  "body":[B("plain","8 degradation classes, each scored 1–5 from its exact applied parameter (5=pristine,1=severe)."),
          B("h3","Three design choices"),
          B("bullet","Perceptual bands — thresholds follow speech-quality convention (telephone BW ≈ 3.4 kHz)."),
          B("bullet","Reverb floors at 4 for any applied RIR — measured RIRs measured at distance, never truly dry;"),
          B("sub","very low DRR (<−15 dB) drops one more band. (Fixed an early RT60-0.15-scored-clean mislabel.)"),
          B("bullet","Use the measured effect, not the knob — clip/disc. on measured fraction; codec→measured rolloff.")],
  "table":{"colw":[0.24,0.76],"cols":["axis","parameter → score bands"],
           "rows":[["noise","SNR: ≥30→5, 20-30→4, 12-20→3, 6-12→2, <6→1"],
                   ["reverb","RT60: <.25→4, .45→3, .70→2, >.70→1; DRR<−15 → −1"],
                   ["bandwidth","cutoff: ≥7500→5, 6000→4, 4000→3, 2500→2, <2500→1"],
                   ["clipping","frac: 0→5, <1%→4, 5%→3, 15%→2, >15%→1"],
                   ["discontinuity","loss: 0→5, <2%→4, 5%→3, 10%→2, >10%→1"],
                   ["loudness","|Δgain|: ≤3→5, 6→4, 12→3, 20→2, >20→1"]]}},

 # 10 targets
 {"title":"Targets: calibrate-then-describe",
  "body":[B("h3","Per-dimension 1–5 scores"),
          B("bullet","From exact params — direct fix for 'quality prior overrides the named dimension'."),
          B("h3","Overall MOS"),
          B("bullet","0.55·min + 0.45·mean of dim scores (worst-axis-dominant), blended 70/30 with fused"),
          B("sub","PESQ+NISQA+DNSMOS. Pure metric fusion rejected: PESQ floors on reverb, DNSMOS reverb-blind."),
          B("h3","Expected I/O"),
          B("bullet","Input: waveform + sqa_full instruction."),
          B("bullet","Stage 1 target: score block ONLY (learn calibrated numbers first)."),
          B("bullet","Stage 2 target: scores + paraphrased description + Overall MOS."),
          B("dim","Cross-entropy on target tokens only; prompt + 88 audio embeds masked to −100.")]},

 # 11 training / LoRA
 {"title":"Training: two-stage LoRA-SFT",
  "body":[B("h3","What LoRA does here"),
          B("bullet","Learn tiny low-rank adapters (r=8, α=28) in Vicuna's projections — a few M params that"),
          B("sub","re-shape the read-out without disturbing the 7B frozen base."),
          B("bullet","Trainable: Q-Former + speech→LLaMA proj + Vicuna-LoRA. Frozen: Whisper, BEATs, Vicuna base."),
          B("dim","single GPU · AMP · batch 2 × accum 8 · lr 1e-5 · warmup→cosine"),
          B("h3","Why two stages"),
          B("bullet","Learn calibrated numbers in isolation first (no description tokens competing for gradient),"),
          B("sub","then describe + emit MOS on top of a representation that already rates correctly.")],
  "table":{"colw":[0.22,0.40,0.26,0.12],"cols":["stage","target text","from","ep"],
           "rows":[["1 — calibration\n(sqa_score)","per-dimension score block only","released SQA ckpt","2"],
                   ["2 — reasoning\n(sqa_full)","scores + description + Overall MOS","Stage 1 best","4"]]}},

 # 12 results table
 {"title":"Results: all five findings re-measured",
  "body":[B("dim","Held-out synthetic degradation sweep on VoiceBank-DEMAND clean (disjoint speakers).")],
  "table":{"colw":[0.05,0.40,0.55],"cols":["#","Original finding","v3"],
           "rows":[["1","MOS↔SNR ρ 0.37, lenient","ρ 0.50; penalizes low SNR; names noise 99% (low SNR)"],
                   ["2","ρ 0.40–0.49; 5 values floored","ρ 0.69–0.75 (in metrics' band); 81 distinct (2.10–4.89)"],
                   ["3","reverb −0.11, bw −0.39, clip −0.33","reverb −0.95, bw −0.81, clip −0.83, noise −0.93"],
                   ["4","Enhancement MOS +0.03, ρ≈0","MOS gain +0.68 (76% of files), ρ +0.22"],
                   ["5","Acquiescence → always YES","calibrated discriminative scores; N/A"]]},
  "kpis":[("5→81","distinct MOS values"),("−.11→−.95","reverb ρ"),
          ("0 / 108","degenerate outputs"),("90%","real-noise naming")]},

 # 13 results calibration plots
 {"title":"Results: the MOS scale is rebuilt",
  "body":[B("bullet","Left: MOS vs SNR — v3 cloud spreads the full range and slopes with SNR (was a flat band)."),
          B("bullet","Right: v3 MOS vs metrics — continuous, rank-aligned (ρ 0.69–0.75); stripes gone.")],
  "image_pair":[os.path.join(PLB,"mos_vs_snr_v3.png"), os.path.join(PLB,"mos_vs_neural_v3.png")]},

 # 14 sweep + enhancement
 {"title":"Results: every axis now responds",
  "body":[B("bullet","Left: MOS vs severity — monotone downward everywhere, incl. the former reverb blind spot."),
          B("bullet","Right: enhancement — v3 MOS rises with denoiser gain (usable as an enhancer evaluator).")],
  "image_pair":[os.path.join(PLB,"degradation_sweep_v3.png"), os.path.join(PLB,"enhancement_blindness.png")]},

 # 15 scale usage
 {"title":"Results: from a 5-rung ladder to a continuous meter",
  "body":[B("bullet","Original emitted a handful of MOS values, floored well above the bottom of the scale."),
          B("good","v3 uses 81 distinct values spanning 2.10–4.89 — it meters quality instead of bucketing it."),
          B("dim","This de-compression lets rank-correlations climb into the band where the metrics agree.")],
  "image":os.path.join(PLB,"mos_scale_usage.png")},

 # 16 bottom line
 {"title":"Bottom line",
  "body":[B("plain","v3 turns a lenient noise specialist with a weak 5-value MOS into a calibrated rater:"),
          B("bullet","MOS de-compressed (5→81), agreeing with PESQ/NISQA/DNSMOS as well as they agree with each other"),
          B("bullet","Perceives and ranks every degradation type — incl. the former reverb blind spot"),
          B("bullet","Tracks enhancement gain → usable as an automatic enhancer evaluator (it wasn't before)"),
          B("bullet","Names degradations in natural, grounded prose; 0 degenerate outputs"),
          B("dim","Per-dimension structured scores are the most reliable signal. The whole fix was data + targets —"),
          B("dim","encoders never unfrozen, validated up-front by the Stage 0 probe."),
          B("h3","Optional further work"),
          B("bullet","Scale clean pool to train-clean-360 · more degradation types · scalar regression head if banding remains.")]},
]

# ======================================================================
# PPTX backend
# ======================================================================
EMU = 914400
SW, SH = 13.333, 7.5
def _fs(slide, x, y, w, h):  # full-slide rect helper
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def build_pptx(path):
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]
    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        # background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = H(BG); bg.line.fill.background()
        bg.shadow.inherit = False
        slide.shapes._spTree.remove(bg._element); slide.shapes._spTree.insert(2, bg._element)

        is_title = s.get("type") == "title"
        # title
        if is_title:
            tb = _fs(slide, 0.8, 2.3, 11.7, 2.2); tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = s["title"]
            r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = H(INK)
            p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = s["subtitle"]
            r2.font.size = Pt(22); r2.font.color.rgb = H(DIM)
            body_top = 4.7
        else:
            tb = _fs(slide, 0.55, 0.3, 12.2, 1.0); tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = s["title"]
            r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = H(ACC)
            # accent underline
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.18), Inches(12.1), Pt(2))
            ln.fill.solid(); ln.fill.fore_color.rgb = H(LINE); ln.line.fill.background(); ln.shadow.inherit = False
            body_top = 1.45

        has_img = "image" in s
        has_pair = "image_pair" in s
        has_tbl = "table" in s
        has_top = "topimage" in s
        # body text region width depends on layout (tables now stack full-width below body)
        if has_img:
            txt_w = 6.6
        else:
            txt_w = 12.2
        ty = body_top + (1.6 if has_top else 0)

        # top image (pipeline)
        if has_top:
            slide.shapes.add_picture(s["topimage"], Inches(0.7), Inches(body_top), width=Inches(11.9))

        # body lines
        if s.get("body"):
            tb = _fs(slide, 0.65, ty, txt_w, SH - ty - 0.3); tf = tb.text_frame; tf.word_wrap = True
            first = True
            for kind, text in s["body"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                p.space_after = Pt(4)
                run = p.add_run(); run.text = ("•  " if kind == "bullet" else
                                               "    – " if kind == "sub" else "") + text
                fc = KIND_COLOR.get(kind, INK)
                run.font.color.rgb = H(fc)
                if kind == "h3":
                    run.font.size = Pt(17 if not is_title else 18); run.font.bold = True
                elif kind in ("good","bad","warn"):
                    run.font.size = Pt(15); run.font.bold = True
                elif kind == "sub":
                    run.font.size = Pt(13)
                elif kind == "dim":
                    run.font.size = Pt(13)
                elif is_title:
                    run.font.size = Pt(18)
                else:
                    run.font.size = Pt(14.5)

        # single image (right)
        if has_img:
            from PIL import Image
            iw, ih = Image.open(s["image"]).size
            maxw, maxh = 6.0, SH - body_top - 0.5
            scale = min(maxw/(iw/96.0), maxh/(ih/96.0))
            w_in = (iw/96.0)*scale
            slide.shapes.add_picture(s["image"], Inches(7.0), Inches(body_top+0.1), width=Inches(w_in))

        # image pair (two side by side, large)
        if has_pair:
            from PIL import Image
            n = len(s["image_pair"]); top = body_top + 0.9
            cellw = (12.2 / n)
            for i, img in enumerate(s["image_pair"]):
                iw, ih = Image.open(img).size
                maxw, maxh = cellw - 0.3, SH - top - 0.4
                scale = min(maxw/(iw/96.0), maxh/(ih/96.0))
                w_in = (iw/96.0)*scale; h_in = (ih/96.0)*scale
                x = 0.6 + i*cellw + (cellw - 0.3 - w_in)/2
                slide.shapes.add_picture(img, Inches(x), Inches(top + (maxh-h_in)/2), width=Inches(w_in))

        # table: full-width, stacked below the body
        if has_tbl:
            t = s["table"]; cols = t["cols"]; rows = t["rows"]
            nr, nc = len(rows)+1, len(cols)
            tx, tw = 0.6, 12.1
            ttop = body_top + 0.34*len(s.get("body") or []) + 0.2
            bottom_limit = (SH - 1.55) if s.get("kpis") else (SH - 0.35)
            th = min(bottom_limit - ttop, 0.46*nr + 0.2)
            gfx = slide.shapes.add_table(nr, nc, Inches(tx), Inches(ttop), Inches(tw), Inches(th))
            tbl = gfx.table
            if t.get("colw"):
                for ci, frac in enumerate(t["colw"]):
                    tbl.columns[ci].width = Inches(tw*frac)
            # column widths: first/last narrow for numeric where applicable
            for ci, c in enumerate(cols):
                cell = tbl.cell(0, ci); cell.text = c
                para = cell.text_frame.paragraphs[0]; para.runs[0].font.size = Pt(12)
                para.runs[0].font.bold = True; para.runs[0].font.color.rgb = H(ACC)
                cell.fill.solid(); cell.fill.fore_color.rgb = H("#1f2630")
            for ri, row in enumerate(rows, start=1):
                for ci, val in enumerate(row):
                    cell = tbl.cell(ri, ci); cell.text = str(val)
                    pr = cell.text_frame.paragraphs[0].runs
                    if pr:
                        pr[0].font.size = Pt(11); pr[0].font.color.rgb = H(INK)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = H(PANEL if ri % 2 else "#12171e")
            # shrink row heights
            for r_ in tbl.rows:
                r_.height = Inches(th/nr)

        # kpis row
        if s.get("kpis"):
            kp = s["kpis"]; n = len(kp); gap = 0.25
            total = 12.2; w = (total - gap*(n-1))/n; y = SH - 1.35
            for i,(num,lbl) in enumerate(kp):
                x = 0.6 + i*(w+gap)
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.05))
                box.fill.solid(); box.fill.fore_color.rgb = H(PANEL)
                box.line.color.rgb = H(LINE); box.line.width = Pt(1); box.shadow.inherit = False
                tf = box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = num; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = H(GOOD)
                p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run(); r2.text = lbl; r2.font.size = Pt(11); r2.font.color.rgb = H(DIM)

        # footnote
        if s.get("footnote"):
            tb = _fs(slide, 0.65, SH-0.85, 12.0, 0.7); tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = s["footnote"]
            r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = H(DIM)

    prs.save(path)
    return path

# ======================================================================
# PDF backend (matplotlib PdfPages)
# ======================================================================
def _wrap(text, width):
    return textwrap.wrap(text, width=width) or [""]

def build_pdf(path):
    with PdfPages(path) as pdf:
        for s in SLIDES:
            fig = plt.figure(figsize=(SW, SH), facecolor=BG)
            fig.patch.set_facecolor(BG)
            ax = fig.add_axes([0,0,1,1]); ax.set_facecolor(BG); ax.axis("off")
            ax.set_xlim(0,1); ax.set_ylim(0,1)
            is_title = s.get("type") == "title"

            if is_title:
                ax.text(0.06, 0.66, s["title"], color=INK, fontsize=46, weight="bold", va="center")
                for i, line in enumerate(_wrap(s["subtitle"], 70)):
                    ax.text(0.06, 0.55-0.06*i, line, color=DIM, fontsize=19, va="center")
                y = 0.34
                for kind, t in s["body"]:
                    ax.text(0.06, y, t, color=KIND_COLOR.get(kind,INK), fontsize=14, va="center")
                    y -= 0.055
                pdf.savefig(fig, facecolor=BG); plt.close(fig); continue

            # title + underline
            ax.text(0.045, 0.93, s["title"], color=ACC, fontsize=25, weight="bold", va="center")
            ax.plot([0.045,0.955],[0.885,0.885], color=LINE, lw=1)

            has_img = "image" in s; has_pair = "image_pair" in s
            has_tbl = "table" in s; has_top = "topimage" in s
            txt_w = 54 if has_img else 110
            ytop = 0.83
            if has_top:
                im = mpimg.imread(s["topimage"])
                axi = fig.add_axes([0.06, 0.6, 0.88, 0.22]); axi.axis("off"); axi.imshow(im)
                ytop = 0.55

            # body
            y = ytop
            for kind, t in (s.get("body") or []):
                fs = {"h3":16,"good":14,"bad":14,"warn":14,"sub":12,"dim":12}.get(kind,13.5)
                bold = kind in ("h3","good","bad","warn")
                prefix = "•  " if kind=="bullet" else ("     – " if kind=="sub" else "")
                ww = max(20, int(txt_w * 13.5 / fs))
                lines = _wrap(prefix+t, ww)
                for j, ln in enumerate(lines):
                    ax.text(0.05, y, ("" if j==0 else "    ")+ln, color=KIND_COLOR.get(kind,INK),
                            fontsize=fs, weight=("bold" if bold else "normal"), va="top")
                    y -= 0.045
                y -= 0.012

            # single image right
            if has_img:
                im = mpimg.imread(s["image"]); h,w = im.shape[0], im.shape[1]
                aspect = w/h
                boxw, boxh = 0.42, 0.62; bx, by = 0.53, 0.12
                if aspect > boxw/boxh*(SH/SW):
                    pass
                axi = fig.add_axes([0.53, 0.10, 0.44, 0.70]); axi.axis("off"); axi.imshow(im)

            # pair
            if has_pair:
                n = len(s["image_pair"])
                for i, img in enumerate(s["image_pair"]):
                    im = mpimg.imread(img)
                    axi = fig.add_axes([0.04 + i*0.49, 0.10, 0.45, 0.66]); axi.axis("off"); axi.imshow(im)

            # table: full-width, stacked below the body
            if has_tbl:
                t = s["table"]; cols = t["cols"]; rows = t["rows"]
                tbot = 0.20 if s.get("kpis") else 0.05
                ttop = max(y - 0.02, tbot + 0.10)
                axi = fig.add_axes([0.05, tbot, 0.90, ttop - tbot]); axi.axis("off")
                cell_text = [[c.replace("\n"," / ") for c in row] for row in rows]
                tbl = axi.table(cellText=cell_text, colLabels=cols, loc="center", cellLoc="left",
                                colWidths=t.get("colw"))
                tbl.auto_set_font_size(False); tbl.set_fontsize(10); tbl.scale(1, 1.4)
                for (ri, ci), cell in tbl.get_celld().items():
                    cell.set_edgecolor(LINE)
                    if ri == 0:
                        cell.set_facecolor("#1f2630"); cell.get_text().set_color(ACC); cell.get_text().set_weight("bold")
                    else:
                        cell.set_facecolor(PANEL if ri % 2 else "#12171e"); cell.get_text().set_color(INK)

            # kpis
            if s.get("kpis"):
                kp = s["kpis"]; n=len(kp); gap=0.015; total=0.9; w=(total-gap*(n-1))/n
                for i,(num,lbl) in enumerate(kp):
                    x = 0.05 + i*(w+gap)
                    ax.add_patch(FancyBboxPatch((x,0.04), w, 0.13, boxstyle="round,pad=0.005,rounding_size=0.01",
                                 transform=ax.transAxes, linewidth=1, edgecolor=LINE, facecolor=PANEL))
                    ax.text(x+w/2, 0.13, num, ha="center", color=GOOD, fontsize=19, weight="bold")
                    ax.text(x+w/2, 0.07, lbl, ha="center", color=DIM, fontsize=10)

            if s.get("footnote"):
                for j, ln in enumerate(_wrap(s["footnote"], 130)):
                    ax.text(0.05, 0.05-0.03*j, ln, color=DIM, fontsize=10.5, style="italic", va="center")

            pdf.savefig(fig, facecolor=BG); plt.close(fig)
    return path

if __name__ == "__main__":
    p1 = build_pptx(os.path.join(RES, "planb_deck.pptx"))
    print("wrote", p1)
    p2 = build_pdf(os.path.join(RES, "planb_deck.pdf"))
    print("wrote", p2)
